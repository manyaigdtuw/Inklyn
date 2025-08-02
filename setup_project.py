#!/usr/bin/env python3
"""
Automatic Project Setup Script for AI Document Chatbot with OpenRouter
"""

import os
import sys

def create_directory_structure():
    """Create all necessary directories"""
    directories = [
        'utils',
        'static/uploads',
        'templates',
        'temp_files'
    ]
    
    for directory in directories:
        os.makedirs(directory, exist_ok=True)
        print(f"✅ Created directory: {directory}")
    
    # Create __init__.py files for Python packages
    init_files = ['utils/__init__.py']
    for init_file in init_files:
        with open(init_file, 'w') as f:
            f.write('# Init file\n')
        print(f"✅ Created: {init_file}")

def create_requirements_txt():
    """Create requirements.txt file"""
    requirements = """streamlit==1.29.0
openai==1.3.8
python-dotenv==1.0.0
PyPDF2==3.0.1
python-docx==1.1.0
pandas==2.1.4
Pillow==10.1.0
pytesseract==0.3.10
opencv-python==4.8.1.78
easyocr==1.7.0
pymupdf==1.23.9
openpyxl==3.1.2
langdetect==1.0.9
streamlit-chat==0.1.1
streamlit-option-menu==0.3.6
extra-streamlit-components==0.1.60
requests==2.31.0
httpx==0.25.2
"""
    
    with open('requirements.txt', 'w') as f:
        f.write(requirements)
    print("✅ Created: requirements.txt")

def create_env_file():
    """Create .env template file"""
    env_content = """# OpenRouter API Configuration
OPENROUTER_API_KEY=your_openrouter_api_key_here
OPENROUTER_BASE_URL=https://openrouter.ai/api/v1

# Optional: Your app name for OpenRouter
APP_NAME=AI Document Chatbot

# Model Selection (you can change this)
DEFAULT_MODEL=anthropic/claude-3-haiku

# Alternative models you can use:
# openai/gpt-4-turbo-preview
# openai/gpt-3.5-turbo
# anthropic/claude-3-sonnet
# meta-llama/llama-2-70b-chat
# google/palm-2-chat-bison
"""
    
    with open('.env', 'w') as f:
        f.write(env_content)
    print("✅ Created: .env (Please add your OpenRouter API key!)")

def create_config_py():
    """Create config.py file"""
    config_content = '''import os
from dotenv import load_dotenv

load_dotenv()

class Config:
    # OpenRouter Configuration
    OPENROUTER_API_KEY = os.environ.get('OPENROUTER_API_KEY')
    OPENROUTER_BASE_URL = os.environ.get('OPENROUTER_BASE_URL', 'https://openrouter.ai/api/v1')
    APP_NAME = os.environ.get('APP_NAME', 'AI Document Chatbot')
    
    # Model settings
    DEFAULT_MODEL = os.environ.get('DEFAULT_MODEL', 'anthropic/claude-3-haiku')
    MAX_TOKENS = 2000
    TEMPERATURE = 0.7
    
    # File settings
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    SUPPORTED_FORMATS = [
        'pdf', 'docx', 'doc', 'txt', 'csv', 'xlsx', 'xls',
        'png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'webp',
        'pptx', 'json', 'html', 'md', 'rtf'
    ]
    
    # UI settings
    PAGE_TITLE = "AI Document Assistant"
    PAGE_ICON = "🤖"
    
    # Available models (you can add more)
    AVAILABLE_MODELS = {
        'Claude 3 Haiku': 'anthropic/claude-3-haiku',
        'Claude 3 Sonnet': 'anthropic/claude-3-sonnet',
        'GPT-4 Turbo': 'openai/gpt-4-turbo-preview',
        'GPT-3.5 Turbo': 'openai/gpt-3.5-turbo',
        'Llama 2 70B': 'meta-llama/llama-2-70b-chat',
        'Gemini Pro': 'google/gemini-pro'
    }
'''
    
    with open('config.py', 'w') as f:
        f.write(config_content)
    print("✅ Created: config.py")

def create_document_processor():
    """Create document_processor.py"""
    processor_content = '''import PyPDF2
import docx
import pandas as pd
from PIL import Image
import pytesseract
import easyocr
import fitz  # PyMuPDF
import json
import io
import cv2
import numpy as np
from pptx import Presentation
from config import Config
import streamlit as st

class DocumentProcessor:
    def __init__(self):
        try:
            self.ocr_reader = easyocr.Reader(['en'])
        except Exception:
            self.ocr_reader = None
    
    def process_file(self, file_content, filename, file_type):
        """Main file processing method"""
        try:
            file_ext = filename.split('.')[-1].lower()
            
            # Route to appropriate processor
            if file_ext == 'pdf':
                return self._process_pdf(file_content)
            elif file_ext in ['docx', 'doc']:
                return self._process_word(file_content)
            elif file_ext == 'txt':
                return self._process_text(file_content)
            elif file_ext in ['csv']:
                return self._process_csv(file_content)
            elif file_ext in ['xlsx', 'xls']:
                return self._process_excel(file_content)
            elif file_ext in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'webp']:
                return self._process_image(file_content)
            elif file_ext == 'pptx':
                return self._process_powerpoint(file_content)
            elif file_ext == 'json':
                return self._process_json(file_content)
            else:
                return self._process_generic(file_content)
                
        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'content': '',
                'metadata': {}
            }
    
    def _process_pdf(self, file_content):
        """Process PDF files"""
        try:
            content = ""
            metadata = {}
            
            # Try PyMuPDF first
            doc = fitz.open(stream=file_content, filetype="pdf")
            
            for page_num in range(doc.page_count):
                page = doc[page_num]
                content += page.get_text() + "\\n"
            
            metadata = {
                'pages': doc.page_count,
                'title': doc.metadata.get('title', 'Unknown')
            }
            
            doc.close()
            
            return {
                'success': True,
                'content': content.strip(),
                'metadata': metadata,
                'type': 'PDF Document'
            }
            
        except Exception as e:
            return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}
    
    def _process_word(self, file_content):
        """Process Word documents"""
        try:
            doc = docx.Document(io.BytesIO(file_content))
            content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)
            
            # Extract tables
            tables_content = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(' | '.join(row_data))
                tables_content.append('\\n'.join(table_data))
            
            full_content = '\\n'.join(content)
            if tables_content:
                full_content += '\\n\\nTables:\\n' + '\\n\\n'.join(tables_content)
            
            return {
                'success': True,
                'content': full_content,
                'metadata': {
                    'paragraphs': len(content),
                    'tables': len(tables_content)
                },
                'type': 'Word Document'
            }
            
        except Exception as e:
            return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}
    
    def _process_text(self, file_content):
        """Process text files"""
        try:
            content = file_content.decode('utf-8')
            return {
                'success': True,
                'content': content,
                'metadata': {
                    'lines': len(content.split('\\n')),
                    'characters': len(content)
                },
                'type': 'Text File'
            }
        except UnicodeDecodeError:
            try:
                content = file_content.decode('latin-1')
                return {
                    'success': True,
                    'content': content,
                    'metadata': {'encoding': 'latin-1'},
                    'type': 'Text File'
                }
            except Exception as e:
                return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}
    
    def _process_csv(self, file_content):
        """Process CSV files"""
        try:
            df = pd.read_csv(io.BytesIO(file_content))
            
            content = f"CSV Data Analysis:\\n"
            content += f"Rows: {len(df)}, Columns: {len(df.columns)}\\n"
            content += f"Columns: {', '.join(df.columns.tolist())}\\n\\n"
            content += "First 10 rows:\\n"
            content += df.head(10).to_string()
            
            return {
                'success': True,
                'content': content,
                'metadata': {
                    'rows': len(df),
                    'columns': len(df.columns),
                    'column_names': df.columns.tolist()
                },
                'type': 'CSV File'
            }
            
        except Exception as e:
            return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}
    
    def _process_excel(self, file_content):
        """Process Excel files"""
        try:
            df = pd.read_excel(io.BytesIO(file_content))
            
            content = f"Excel Data Analysis:\\n"
            content += f"Rows: {len(df)}, Columns: {len(df.columns)}\\n"
            content += f"Columns: {', '.join(df.columns.tolist())}\\n\\n"
            content += "First 10 rows:\\n"
            content += df.head(10).to_string()
            
            return {
                'success': True,
                'content': content,
                'metadata': {
                    'rows': len(df),
                    'columns': len(df.columns),
                    'column_names': df.columns.tolist()
                },
                'type': 'Excel File'
            }
            
        except Exception as e:
            return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}
    
    def _process_image(self, file_content):
        """Process image files with OCR"""
        try:
            # Convert bytes to image
            image = Image.open(io.BytesIO(file_content))
            
            # Try multiple OCR methods
            text_content = ""
            
            # Method 1: pytesseract
            try:
                text_content = pytesseract.image_to_string(image)
            except:
                pass
            
            # Method 2: EasyOCR (if available)
            if not text_content.strip() and self.ocr_reader:
                try:
                    # Convert PIL to numpy array for EasyOCR
                    img_array = np.array(image)
                    results = self.ocr_reader.readtext(img_array)
                    text_content = ' '.join([result[1] for result in results])
                except:
                    pass
            
            if not text_content.strip():
                text_content = "No text detected in image"
            
            return {
                'success': True,
                'content': f"Image OCR Results:\\n{text_content}",
                'metadata': {
                    'size': image.size,
                    'mode': image.mode,
                    'format': image.format
                },
                'type': 'Image File'
            }
            
        except Exception as e:
            return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}
    
    def _process_powerpoint(self, file_content):
        """Process PowerPoint files"""
        try:
            prs = Presentation(io.BytesIO(file_content))
            content = []
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = f"Slide {slide_num}:\\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content += f"- {shape.text.strip()}\\n"
                content.append(slide_content)
            
            full_content = "\\n".join(content)
            
            return {
                'success': True,
                'content': full_content,
                'metadata': {
                    'slides': len(prs.slides)
                },
                'type': 'PowerPoint Presentation'
            }
            
        except Exception as e:
            return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}
    
    def _process_json(self, file_content):
        """Process JSON files"""
        try:
            data = json.loads(file_content.decode('utf-8'))
            
            # Pretty print JSON
            content = f"JSON File Analysis:\\n"
            content += json.dumps(data, indent=2, ensure_ascii=False)
            
            return {
                'success': True,
                'content': content,
                'metadata': {
                    'type': type(data).__name__,
                    'keys': list(data.keys()) if isinstance(data, dict) else None
                },
                'type': 'JSON File'
            }
            
        except Exception as e:
            return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}
    
    def _process_generic(self, file_content):
        """Generic text processing for unknown file types"""
        try:
            content = file_content.decode('utf-8')
            return {
                'success': True,
                'content': content,
                'metadata': {'size': len(content)},
                'type': 'Generic Text'
            }
        except:
            return {
                'success': False,
                'error': 'Cannot process this file type',
                'content': '',
                'metadata': {}
            }
'''
    
    with open('document_processor.py', 'w') as f:
        f.write(processor_content)
    print("✅ Created: document_processor.py")

def create_ai_engine():
    """Create ai_engine.py with OpenRouter integration"""
    ai_engine_content = '''import httpx
import json
from config import Config
import streamlit as st

class AIEngine:
    def __init__(self):
        self.api_key = Config.OPENROUTER_API_KEY
        self.base_url = Config.OPENROUTER_BASE_URL
        self.app_name = Config.APP_NAME
        self.model = Config.DEFAULT_MODEL
        
        if not self.api_key:
            st.error("⚠️ OpenRouter API key not found! Please add it to your .env file")
    
    def generate_response(self, user_message, context, chat_history):
        """Generate AI response using OpenRouter"""
        try:
            # Build conversation context
            messages = self._build_messages(user_message, context, chat_history)
            
            # Make API call to OpenRouter
            response = self._call_openrouter_api(messages)
            
            return response
            
        except Exception as e:
            return f"I apologize, but I encountered an error: {str(e)}. Please try again."
    
    def _build_messages(self, user_message, context, chat_history):
        """Build message array for API call"""
        messages = []
        
        # System prompt
        system_prompt = f"""You are an AI Document Assistant that helps users understand documents, write emails, and answer questions.

Current Context from uploaded files:
{context}

You can help with:
1. 📖 Explaining document content
2. 📝 Writing professional emails
3. 📧 Replying to emails with context
4. 📊 Analyzing data and documents  
5. ❓ Answering questions about uploaded files

Be helpful, professional, and use the document context when relevant."""

        messages.append({"role": "system", "content": system_prompt})
        
        # Add recent chat history (last 10 messages)
        recent_history = chat_history[-10:] if len(chat_history) > 10 else chat_history
        for msg in recent_history:
            if msg['role'] in ['user', 'assistant']:
                messages.append({
                    "role": msg['role'],
                    "content": msg['content']
                })
        
        return messages
    
    def _call_openrouter_api(self, messages):
        """Make API call to OpenRouter"""
        headers = {
            "Authorization": f"Bearer {self.api_key}",
            "HTTP-Referer": "http://localhost:8501",  # Streamlit default
            "X-Title": self.app_name,
            "Content-Type": "application/json"
        }
        
        data = {
            "model": self.model,
            "messages": messages,
            "max_tokens": Config.MAX_TOKENS,
            "temperature": Config.TEMPERATURE,
            "stream": False
        }
        
        with httpx.Client(timeout=30.0) as client:
            response = client.post(
                f"{self.base_url}/chat/completions",
                headers=headers,
                json=data
            )
            
            if response.status_code == 200:
                result = response.json()
                return result['choices'][0]['message']['content'].strip()
            else:
                error_msg = f"API Error {response.status_code}: {response.text}"
                st.error(error_msg)
                return "I'm having trouble connecting to the AI service. Please try again."
    
    def change_model(self, model_name):
        """Change the AI model"""
        if model_name in Config.AVAILABLE_MODELS.values():
            self.model = model_name
            return True
        return False
    
    def get_available_models(self):
        """Get list of available models"""
        return Config.AVAILABLE_MODELS
'''
    
    with open('ai_engine.py', 'w') as f:
        f.write(ai_engine_content)
    print("✅ Created: ai_engine.py")

def create_email_assistant():
    """Create email_assistant.py"""
    email_content = '''from ai_engine import AIEngine

class EmailAssistant:
    def __init__(self):
        self.ai_engine = AIEngine()
    
    def write_email(self, context, requirements):
        """Help write an email with document context"""
        prompt = f"""
Please help me write a professional email based on the following:

Requirements: {requirements}

Document Context: {context}

Please provide:
1. Subject line
2. Professional email body
3. Appropriate greeting and closing

Make it clear, professional, and well-structured.
"""
        return self.ai_engine._call_openrouter_api([
            {"role": "system", "content": "You are a professional email writing assistant."},
            {"role": "user", "content": prompt}
        ])
    
    def reply_to_email(self, original_email, context, instructions=""):
        """Help reply to an email"""
        prompt = f"""
Please help me write a professional reply to this email:

Original Email:
{original_email}

Context from my documents: {context}

Additional instructions: {instructions}

Please write an appropriate professional reply.
"""
        return self.ai_engine._call_openrouter_api([
            {"role": "system", "content": "You are a professional email reply assistant."},
            {"role": "user", "content": prompt}
        ])
'''
    
    with open('email_assistant.py', 'w') as f:
        f.write(email_content)
    print("✅ Created: email_assistant.py")

def create_main_app():
    """Create the main Streamlit app.py"""
    app_content = '''import streamlit as st
import openai
from datetime import datetime
import json
import uuid
from document_processor import DocumentProcessor
from ai_engine import AIEngine
from email_assistant import EmailAssistant
from config import Config
from streamlit_option_menu import option_menu
import extra_streamlit_components as stx

# Configure page
st.set_page_config(
    page_title=Config.PAGE_TITLE,
    page_icon=Config.PAGE_ICON,
    layout="wide",
    initial_sidebar_state="expanded"
)

# Initialize components
@st.cache_resource
def init_components():
    return {
        'doc_processor': DocumentProcessor(),
        'ai_engine': AIEngine(),
        'email_assistant': EmailAssistant()
    }

components = init_components()

# Initialize session state
def init_session_state():
    if 'session_id' not in st.session_state:
        st.session_state.session_id = str(uuid.uuid4())
    if 'chat_history' not in st.session_state:
        st.session_state.chat_history = []
    if 'uploaded_files' not in st.session_state:
        st.session_state.uploaded_files = []
    if 'file_contents' not in st.session_state:
        st.session_state.file_contents = {}
    if 'current_context' not in st.session_state:
        st.session_state.current_context = ""
    if 'selected_model' not in st.session_state:
        st.session_state.selected_model = Config.DEFAULT_MODEL

init_session_state()

# Custom CSS
def load_css():
    st.markdown("""
    <style>
    .main {
        padding-top: 1rem;
    }
    
    .chat-message {
        padding: 1rem;
        border-radius: 10px;
        margin-bottom: 1rem;
        box-shadow: 0 2px 4px rgba(0,0,0,0.1);
    }
    
    .user-message {
        background-color: #e3f2fd;
        border-left: 4px solid #2196f3;
    }
    
    .assistant-message {
        background-color: #f5f5f5;
        border-left: 4px solid #4caf50;
    }
    
    .file-info {
        background-color: #fff3e0;
        padding: 10px;
        border-radius: 5px;
        border-left: 4px solid #ff9800;
        margin: 10px 0;
    }
    
    .success-message {
        background-color: #e8f5e8;
        color: #2e7d32;
        padding: 10px;
        border-radius: 5px;
        border-left: 4px solid #4caf50;
    }
    
    .error-message {
        background-color: #ffebee;
        color: #c62828;
        padding: 10px;
        border-radius: 5px;
        border-left: 4px solid #f44336;
    }
    
    .stButton > button {
        width: 100%;
        border-radius: 20px;
        border: none;
        background: linear-gradient(90deg, #667eea 0%, #764ba2 100%);
        color: white;
    }
    
    .stTextInput > div > div > input {
        border-radius: 20px;
    }
    
    .model-selector {
        padding: 10px;
        background: #f0f0f0;
        border-radius: 10px;
        margin-bottom: 20px;
    }
    </style>
    """, unsafe_allow_html=True)

load_css()

# Main App Layout
def main():
    st.title("🤖 AI Document Assistant (OpenRouter)")
    st.markdown("**Upload any file and I'll help you understand it, write emails, or answer questions!**")
    
    # Sidebar for file management and settings
    with st.sidebar:
        st.header("⚙️ Settings")
        
        # Model selection
        st.markdown('<div class="model-selector">', unsafe_allow_html=True)
        st.subheader("🧠 AI Model")
        available_models = components['ai_engine'].get_available_models()
        
        selected_model_name = st.selectbox(
            "Choose AI Model:",
            list(available_models.keys()),
            index=list(available_models.values()).index(st.session_state.selected_model)
        )
        
        if available_models[selected_model_name] != st.session_state.selected_model:
            st.session_state.selected_model = available_models[selected_model_name]
            components['ai_engine'].change_model(st.session_state.selected_model)
            st.success(f"Switched to {selected_model_name}")
        
        st.markdown('</div>', unsafe_allow_html=True)
        
        st.header("📁 File Manager")
        
        # File upload
        uploaded_files = st.file_uploader(
            "Upload Documents",
            type=Config.SUPPORTED_FORMATS,
            accept_multiple_files=True,
            help="Supported: PDF, Word, Excel, Images, PowerPoint, and more!"
        )
        
        # Process uploaded files
        if uploaded_files:
            process_uploaded_files(uploaded_files)
        
        # Display current files
        display_current_files()
        
        # Quick actions
        st.subheader("⚡ Quick Actions")
        if st.button("📝 Help me write an email"):
            st.session_state.quick_action = "write_email"
        if st.button("📧 Help me reply to an email"):
            st.session_state.quick_action = "reply_email"
        if st.button("📊 Analyze my documents"):
            st.session_state.quick_action = "analyze_docs"
        if st.button("🔄 Clear conversation"):
            clear_conversation()
    
    # Main chat interface
    col1, col2 = st.columns([3, 1])
    
    with col1:
        # Chat history
        display_chat_history()
        
        # Chat input
        handle_chat_input()
    
    with col2:
        # Context panel
        display_context_panel()

def process_uploaded_files(uploaded_files):
    """Process and store uploaded files"""
    for uploaded_file in uploaded_files:
        if uploaded_file.name not in [f['name'] for f in st.session_state.uploaded_files]:
            with st.spinner(f"Processing {uploaded_file.name}..."):
                try:
                    # Save file temporarily
                    file_content = uploaded_file.read()
                    
                    # Process with document processor
                    result = components['doc_processor'].process_file(
                        file_content, uploaded_file.name, uploaded_file.type
                    )
                    
                    if result['success']:
                        # Store file info
                        file_info = {
                            'name': uploaded_file.name,
                            'type': uploaded_file.type,
                            'size': len(file_content),
                            'content': result['content'],
                            'metadata': result['metadata'],
                            'processed_at': datetime.now().strftime("%Y-%m-%d %H:%M")
                        }
                        
                        st.session_state.uploaded_files.append(file_info)
                        st.session_state.file_contents[uploaded_file.name] = result['content']
                        
                        # Update context
                        update_context()
                        
                        st.success(f"✅ Successfully processed {uploaded_file.name}")
                    else:
                        st.error(f"❌ Error processing {uploaded_file.name}: {result['error']}")
                        
                except Exception as e:
                    st.error(f"❌ Error processing {uploaded_file.name}: {str(e)}")

def display_current_files():
    """Display currently uploaded files"""
    if st.session_state.uploaded_files:
        st.subheader("📋 Current Files")
        for i, file_info in enumerate(st.session_state.uploaded_files):
            with st.expander(f"📄 {file_info['name']}", expanded=False):
                st.markdown(f"**Type:** {file_info['type']}")
                st.markdown(f"**Size:** {file_info['size']:,} bytes")
                st.markdown(f"**Processed:** {file_info['processed_at']}")
                
                # Show content preview
                content_preview = file_info['content'][:300] + "..." if len(file_info['content']) > 300 else file_info['content']
                st.text_area("Content Preview", content_preview, height=100, disabled=True)
                
                # Remove file option
                if st.button(f"🗑️ Remove", key=f"remove_{i}"):
                    st.session_state.uploaded_files.pop(i)
                    del st.session_state.file_contents[file_info['name']]
                    update_context()
                    st.rerun()

def display_chat_history():
    """Display chat conversation"""
    if st.session_state.chat_history:
        for message in st.session_state.chat_history:
            if message['role'] == 'user':
                st.markdown(f"""
                <div class="chat-message user-message">
                    <strong>You:</strong><br>
                    {message['content']}
                </div>
                """, unsafe_allow_html=True)
            else:
                st.markdown(f"""
                <div class="chat-message assistant-message">
                    <strong>AI Assistant:</strong><br>
                    {message['content']}
                </div>
                """, unsafe_allow_html=True)
    else:
        st.markdown(f"""
        <div class="chat-message assistant-message">
            <strong>AI Assistant ({st.session_state.selected_model}):</strong><br>
            Hello! I'm your AI Document Assistant powered by OpenRouter. Upload any files and I can help you:
            <ul>
                <li>📖 Understand and explain document content</li>
                <li>📝 Write professional emails</li>
                <li>📧 Reply to emails with context</li>
                <li>📊 Analyze data and documents</li>
                <li>❓ Answer questions about your files</li>
            </ul>
            How can I assist you today?
        </div>
        """, unsafe_allow_html=True)

def handle_chat_input():
    """Handle user chat input"""
    # Check for quick actions
    if 'quick_action' in st.session_state:
        handle_quick_action(st.session_state.quick_action)
        del st.session_state.quick_action
    
    # Chat input
    user_input = st.chat_input("Type your message here...")
    
    if user_input:
        # Add user message to history
        st.session_state.chat_history.append({
            'role': 'user',
            'content': user_input,
            'timestamp': datetime.now()
        })
        
        # Get AI response
        with st.spinner("🤔 Thinking..."):
            ai_response = components['ai_engine'].generate_response(
                user_input,
                st.session_state.current_context,
                st.session_state.chat_history
            )
        
        # Add AI response to history
        st.session_state.chat_history.append({
            'role': 'assistant',
            'content': ai_response,
            'timestamp': datetime.now()
        })
        
        st.rerun()

def handle_quick_action(action):
    """Handle quick action buttons"""
    if action == "write_email":
        st.session_state.chat_history.append({
            'role': 'user',
            'content': "Help me write a professional email using the context from my uploaded files.",
            'timestamp': datetime.now()
        })
    elif action == "reply_email":
        st.session_state.chat_history.append({
            'role': 'user',
            'content': "Help me reply to an email. I'll provide the original email content.",
            'timestamp': datetime.now()
        })
    elif action == "analyze_docs":
        st.session_state.chat_history.append({
            'role': 'user',
            'content': "Please analyze all my uploaded documents and provide a comprehensive summary.",
            'timestamp': datetime.now()
        })

def display_context_panel():
    """Display context and file information panel"""
    st.subheader("🎯 Current Context")
    
    # Model info
    st.info(f"🧠 Using: {st.session_state.selected_model}")
    
    if st.session_state.uploaded_files:
        st.markdown(f"**Files:** {len(st.session_state.uploaded_files)} documents loaded")
        
        # Context summary
        if st.session_state.current_context:
            with st.expander("📝 Context Summary", expanded=True):
                st.text_area(
                    "Current understanding:",
                    st.session_state.current_context[:500] + "..." if len(st.session_state.current_context) > 500 else st.session_state.current_context,
                    height=200,
                    disabled=True
                )
    else:
        st.info("No files uploaded yet. Upload documents to get started!")
    
    # Statistics
    if st.session_state.chat_history:
        st.subheader("📊 Session Stats")
        st.metric("Messages", len(st.session_state.chat_history))
        st.metric("Files Processed", len(st.session_state.uploaded_files))

def update_context():
    """Update current context based on uploaded files"""
    if st.session_state.uploaded_files:
        context_parts = []
        for file_info in st.session_state.uploaded_files:
            context_parts.append(f"File: {file_info['name']}\\nContent: {file_info['content'][:1000]}")
        
        st.session_state.current_context = "\\n\\n".join(context_parts)

def clear_conversation():
    """Clear chat history"""
    st.session_state.chat_history = []
    st.session_state.uploaded_files = []
    st.session_state.file_contents = {}
    st.session_state.current_context = ""
    st.success("🔄 Conversation cleared!")
    st.rerun()

if __name__ == "__main__":
    # Check if OpenRouter API key is configured
    if not Config.OPENROUTER_API_KEY:
        st.error("🚨 OpenRouter API key not configured!")
        st.markdown("""
        Please add your OpenRouter API key to the `.env` file:
        1. Go to [OpenRouter.ai](https://openrouter.ai) and get your API key
        2. Add it to `.env` file: `OPENROUTER_API_KEY=your_key_here`
        3. Restart the application
        """)
        st.stop()
    
    main()
'''
    
    with open('app.py', 'w') as f:
        f.write(app_content)
    print("✅ Created: app.py")

def create_run_script():
    """Create run.py script"""
    run_content = '''#!/usr/bin/env python3
"""
Quick run script for the AI Document Chatbot
"""

import subprocess
import sys
import os

def install_requirements():
    """Install required packages"""
    print("📦 Installing requirements...")
    try:
        subprocess.check_call([sys.executable, "-m", "pip", "install", "-r", "requirements.txt"])
        print("✅ Requirements installed successfully!")
    except subprocess.CalledProcessError as e:
        print(f"❌ Error installing requirements: {e}")
        return False
    return True

def run_app():
    """Run the Streamlit app"""
    print("🚀 Starting AI Document Chatbot...")
    try:
        subprocess.run([sys.executable, "-m", "streamlit", "run", "app.py"])
    except KeyboardInterrupt:
        print("\\n👋 Shutting down gracefully...")
    except Exception as e:
        print(f"❌ Error running app: {e}")

def main():
    print("🤖 AI Document Chatbot Setup")
    print("=" * 40)
    
    # Check if .env exists and has API key
    if not os.path.exists('.env') or 'OPENROUTER_API_KEY=your_openrouter_api_key_here' in open('.env').read():
        print("⚠️  Please configure your OpenRouter API key in .env file first!")
        print("1. Get your API key from https://openrouter.ai")
        print("2. Edit .env file and replace 'your_openrouter_api_key_here' with your actual key")
        print("3. Run this script again")
        return
    
    # Install requirements
    if not install_requirements():
        return
    
    # Run the app
    run_app()

if __name__ == "__main__":
    main()
'''
    
    with open('run.py', 'w') as f:
        f.write(run_content)
    print("✅ Created: run.py")


