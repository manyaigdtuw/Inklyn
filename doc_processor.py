import PyPDF2
import docx
import pandas as pd
from PIL import Image
import pytesseract
import easyocr
import fitz  # PyMuPDF
import json
import io
import cv2
import numpy as np
from pptx import Presentation
from config import Config
import streamlit as st

class DocumentProcessor:
    def __init__(self):
        try:
            self.ocr_reader = easyocr.Reader(['en'])
        except Exception:
            self.ocr_reader = None

    def process_file(self, file_content, filename, file_type):
        """Main file processing method"""
        try:
            file_ext = filename.split('.')[-1].lower()

            if file_ext == 'pdf':
                return self._process_pdf(file_content)
            elif file_ext in ['docx', 'doc']:
                return self._process_word(file_content)
            elif file_ext == 'txt':
                return self._process_text(file_content)
            elif file_ext == 'csv':
                return self._process_csv(file_content)
            elif file_ext in ['xlsx', 'xls']:
                return self._process_excel(file_content)
            elif file_ext in ['png', 'jpg', 'jpeg', 'gif', 'bmp', 'tiff', 'webp']:
                return self._process_image(file_content)
            elif file_ext == 'pptx':
                return self._process_powerpoint(file_content)
            elif file_ext == 'json':
                return self._process_json(file_content)
            else:
                return self._process_generic(file_content)

        except Exception as e:
            return {
                'success': False,
                'error': str(e),
                'content': '',
                'metadata': {}
            }

    def _process_pdf(self, file_content):
        try:
            content = ""
            metadata = {}

            doc = fitz.open(stream=file_content, filetype="pdf")
            for page_num in range(doc.page_count):
                page = doc[page_num]
                content += page.get_text() + "\n"

            metadata = {
                'pages': doc.page_count,
                'title': doc.metadata.get('title', 'Unknown')
            }

            doc.close()

            return {
                'success': True,
                'content': content.strip(),
                'metadata': metadata,
                'type': 'PDF Document'
            }

        except Exception as e:
            return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}

    def _process_word(self, file_content):
        try:
            doc = docx.Document(io.BytesIO(file_content))
            content = []

            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    content.append(paragraph.text)

            tables_content = []
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(' | '.join(row_data))
                tables_content.append('\n'.join(table_data))

            full_content = '\n'.join(content)
            if tables_content:
                full_content += '\n\nTables:\n' + '\n\n'.join(tables_content)

            return {
                'success': True,
                'content': full_content,
                'metadata': {
                    'paragraphs': len(content),
                    'tables': len(tables_content)
                },
                'type': 'Word Document'
            }

        except Exception as e:
            return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}

    def _process_text(self, file_content):
        try:
            content = file_content.decode('utf-8')
            return {
                'success': True,
                'content': content,
                'metadata': {
                    'lines': len(content.split('\n')),
                    'characters': len(content)
                },
                'type': 'Text File'
            }
        except UnicodeDecodeError:
            try:
                content = file_content.decode('latin-1')
                return {
                    'success': True,
                    'content': content,
                    'metadata': {'encoding': 'latin-1'},
                    'type': 'Text File'
                }
            except Exception as e:
                return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}

    def _process_csv(self, file_content):
        try:
            df = pd.read_csv(io.BytesIO(file_content))

            content = f"CSV Data Analysis:\n"
            content += f"Rows: {len(df)}, Columns: {len(df.columns)}\n"
            content += f"Columns: {', '.join(df.columns.tolist())}\n\n"
            content += "First 10 rows:\n"
            content += df.head(10).to_string()

            return {
                'success': True,
                'content': content,
                'metadata': {
                    'rows': len(df),
                    'columns': len(df.columns),
                    'column_names': df.columns.tolist()
                },
                'type': 'CSV File'
            }

        except Exception as e:
            return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}

    def _process_excel(self, file_content):
        try:
            df = pd.read_excel(io.BytesIO(file_content))

            content = f"Excel Data Analysis:\n"
            content += f"Rows: {len(df)}, Columns: {len(df.columns)}\n"
            content += f"Columns: {', '.join(df.columns.tolist())}\n\n"
            content += "First 10 rows:\n"
            content += df.head(10).to_string()

            return {
                'success': True,
                'content': content,
                'metadata': {
                    'rows': len(df),
                    'columns': len(df.columns),
                    'column_names': df.columns.tolist()
                },
                'type': 'Excel File'
            }

        except Exception as e:
            return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}

    def _process_image(self, file_content):
        """Process image files with OCR and log to terminal"""
        try:
            print("\n[INFO] Starting image OCR processing...")

            image = Image.open(io.BytesIO(file_content))
            print(f"[INFO] Image loaded: Format={image.format}, Size={image.size}, Mode={image.mode}")

            gray_image = image.convert('L')
            print("[INFO] Converted image to grayscale for Tesseract")

            text_content = ""

            # Method 1: pytesseract
            try:
                print("[INFO] Running pytesseract...")
                text_content = pytesseract.image_to_string(gray_image)
                if text_content.strip():
                    print("[SUCCESS] Tesseract extracted text:")
                    print("="*40)
                    print(text_content)
                    print("="*40)
                else:
                    print("[WARNING] Tesseract returned no text.")
            except Exception as e:
                print(f"[ERROR] Tesseract OCR failed: {e}")

            # Method 2: EasyOCR (fallback)
            if not text_content.strip() and self.ocr_reader:
                try:
                    print("[INFO] Falling back to EasyOCR...")
                    img_array = np.array(image)
                    results = self.ocr_reader.readtext(img_array)
                    text_content = ' '.join([result[1] for result in results])

                    if text_content.strip():
                        print("[SUCCESS] EasyOCR extracted text:")
                        print("="*40)
                        print(text_content)
                        print("="*40)
                    else:
                        print("[WARNING] EasyOCR also returned no text.")
                except Exception as e:
                    print(f"[ERROR] EasyOCR failed: {e}")

            if not text_content.strip():
                text_content = "No text detected in image"

            return {
                'success': True,
                'content': f"Image OCR Results:\n{text_content}",
                'metadata': {
                    'size': image.size,
                    'mode': image.mode,
                    'format': image.format
                },
                'type': 'Image File'
            }

        except Exception as e:
            print(f"[FATAL] Error in _process_image: {e}")
            return {
                'success': False,
                'error': str(e),
                'content': '',
                'metadata': {}
            }

    def _process_powerpoint(self, file_content):
        try:
            prs = Presentation(io.BytesIO(file_content))
            content = []

            for slide_num, slide in enumerate(prs.slides, 1):
                slide_content = f"Slide {slide_num}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_content += f"- {shape.text.strip()}\n"
                content.append(slide_content)

            full_content = "\n".join(content)

            return {
                'success': True,
                'content': full_content,
                'metadata': {
                    'slides': len(prs.slides)
                },
                'type': 'PowerPoint Presentation'
            }

        except Exception as e:
            return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}

    def _process_json(self, file_content):
        try:
            data = json.loads(file_content.decode('utf-8'))

            content = f"JSON File Analysis:\n"
            content += json.dumps(data, indent=2, ensure_ascii=False)

            return {
                'success': True,
                'content': content,
                'metadata': {
                    'type': type(data).__name__,
                    'keys': list(data.keys()) if isinstance(data, dict) else None
                },
                'type': 'JSON File'
            }

        except Exception as e:
            return {'success': False, 'error': str(e), 'content': '', 'metadata': {}}

    def _process_generic(self, file_content):
        try:
            content = file_content.decode('utf-8')
            return {
                'success': True,
                'content': content,
                'metadata': {'size': len(content)},
                'type': 'Generic Text'
            }
        except:
            return {
                'success': False,
                'error': 'Cannot process this file type',
                'content': '',
                'metadata': {}
            }
